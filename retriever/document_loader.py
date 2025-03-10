from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document as LangchainDocument
from pptx import Presentation
from docx import Document as DocxDocument
import tempfile
import os

def load_and_split_pdf(pdf_bytes):
    """
    Load and split a PDF file into LangchainDocument objects.

    :param pdf_bytes: Bytes of the PDF file.
    :return: List of LangchainDocument objects.
    """
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as temp_file:
            temp_file.write(pdf_bytes.read())
            temp_file_path = temp_file.name

            loader = PyPDFLoader(temp_file_path)
            docs = loader.load()

            # Split documents into chunks
            text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
            return text_splitter.split_documents(docs)

    except Exception as e:
        print(f"❌ Error loading PDF: {e}")
        return []

    finally:
        # Ensure cleanup of temp file
        if 'temp_file_path' in locals():
            os.remove(temp_file_path)


def load_and_split_docx(docx_bytes):
    """
    Load and split a DOCX file into LangchainDocument objects.

    :param docx_bytes: Bytes of the DOCX file.
    :return: List of LangchainDocument objects.
    """
    try:
        doc = DocxDocument(docx_bytes)
        text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])

        # Split text into chunks
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        chunks = text_splitter.split_text(text)

        return [LangchainDocument(page_content=chunk) for chunk in chunks]

    except Exception as e:
        print(f"❌ Error loading DOCX: {e}")
        return []


def load_and_split_pptx(pptx_bytes):
    """
    Load and split a PPTX file into LangchainDocument objects.

    :param pptx_bytes: Bytes of the PPTX file.
    :return: List of LangchainDocument objects.
    """
    try:
        prs = Presentation(pptx_bytes)
        text = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)

        combined_text = "\n".join(text)

        # Split text into chunks
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
        text_chunks = text_splitter.split_text(combined_text)

        return [LangchainDocument(page_content=chunk) for chunk in text_chunks]

    except Exception as e:
        print(f"❌ Error loading PPTX: {e}")
        return []


# Example usage
if __name__ == "__main__":
    # Example with PDF
    with open("example.pdf", "rb") as pdf_file:
        pdf_docs = load_and_split_pdf(pdf_file)
        print(f"Loaded {len(pdf_docs)} PDF documents.")

    # Example with DOCX
    with open("example.docx", "rb") as docx_file:
        docx_docs = load_and_split_docx(docx_file)
        print(f"Loaded {len(docx_docs)} DOCX documents.")

    # Example with PPTX
    with open("example.pptx", "rb") as pptx_file:
        pptx_docs = load_and_split_pptx(pptx_file)
        print(f"Loaded {len(pptx_docs)} PPTX documents.")
