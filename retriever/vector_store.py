from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_core.documents import Document
import os
import streamlit as st
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document as DocxDocument
from concurrent.futures import ThreadPoolExecutor

# Define FAISS database path
VECTOR_DB_PATH = "data/faiss_index"

# Available embedding models
EMBEDDING_MODELS = {
    "minilm": "sentence-transformers/all-MiniLM-L6-v2",
    "bge": "BAAI/bge-base-en",
    "e5": "intfloat/e5-large-v2",  # Using e5-large-v2
    "deepseek": "deepseek-ai/deepseek-embedding"
}

# Select model
MODEL_NAME = EMBEDDING_MODELS["e5"]  # Change to use e5-large-v2


def get_embedding_model(model_name="e5"):
    """Load and return the selected embedding model."""
    model_path = EMBEDDING_MODELS.get(model_name, EMBEDDING_MODELS["e5"])
    return HuggingFaceEmbeddings(model_name=model_path)


# Initialize embeddings once (avoid redundant calls)
embedding = get_embedding_model("e5")


def load_vector_store():
    """Load FAISS vector store if it exists; otherwise, return None."""
    if os.path.exists(VECTOR_DB_PATH):
        try:
            vector_store = FAISS.load_local(VECTOR_DB_PATH, embedding, allow_dangerous_deserialization=True)
            print("✅ FAISS Vector Store Loaded Successfully!")
            return vector_store.as_retriever(search_kwargs={"k": 7})  # Return retriever
        except Exception as e:
            print(f"❌ Error loading FAISS index: {e}")
            return None
    else:
        print("❌ No FAISS Vector Store found!")
        return None


def extract_text_from_file(file):
    """Extract text from PDFs, PPTX, and DOCX files."""
    text = ""

    try:
        if file.name.endswith(".pdf"):
            # Use PyPDF2 for better PDF parsing
            reader = PdfReader(file)
            text = "\n".join(page.extract_text() for page in reader.pages if page.extract_text())

        elif file.name.endswith(".pptx"):
            ppt = Presentation(file)
            extracted_text = [shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text")]
            text = "\n".join(extracted_text)

        elif file.name.endswith(".docx"):
            doc = DocxDocument(file)
            extracted_text = [para.text for para in doc.paragraphs if para.text.strip()]
            text = "\n".join(extracted_text)

        if text.strip():
            print(f"✅ Extracted text from {file.name}: {text[:500]}...")  # Log first 500 characters
        else:
            print(f"⚠️ No text extracted from {file.name}.")

    except Exception as e:
        print(f"❌ Error extracting text from {file.name}: {e}")

    return text


def process_documents(uploaded_files):
    """Process multiple documents in parallel and store them in FAISS."""
    if not uploaded_files:
        print("⚠️ No files uploaded.")
        return None

    docs = []
    with ThreadPoolExecutor() as executor:
        results = list(executor.map(extract_text_from_file, uploaded_files))

    for file, text in zip(uploaded_files, results):
        if text:
            docs.append(Document(page_content=text, metadata={"source": file.name}))

    if not docs:
        print("❌ No valid documents to index.")
        return None

    try:
        if os.path.exists(VECTOR_DB_PATH):
            # Load existing FAISS index and update it
            faiss_index = FAISS.load_local(VECTOR_DB_PATH, embedding, allow_dangerous_deserialization=True)
            faiss_index.add_documents(docs)
            faiss_index.save_local(VECTOR_DB_PATH)
            print("✅ FAISS index updated successfully.")
        else:
            # Create a new FAISS index
            faiss_index = FAISS.from_documents(docs, embedding)
            faiss_index.save_local(VECTOR_DB_PATH)
            print("✅ New FAISS index created.")

        return faiss_index.as_retriever(search_kwargs={"k": 7})  # Return retriever

    except Exception as e:
        print(f"❌ Error updating FAISS index: {e}")
        return None
